from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Configuration
OUTPUT_FILE = "/Users/singhhim/Library/CloudStorage/OneDrive-TheStarsGroup/Documents/himanshu-git/Track-subs/cpitch_deck.ppt"
IMAGES_DIR = "/Users/singhhim/.gemini/antigravity/brain/2ec88451-0d96-41bd-a2d6-0c24d10498f1"

# Image Paths (Update these with actual filenames from generation step)
IMG_TITLE_HERO = os.path.join(IMAGES_DIR, "title_slide_hero_bw_1765140845390.png")
IMG_PROBLEM = os.path.join(IMAGES_DIR, "problem_subscription_chaos_bw_1765140862179.png")
IMG_SOLUTION = os.path.join(IMAGES_DIR, "solution_dashboard_mockup_bw_1765140878653.png")
IMG_MARKET = os.path.join(IMAGES_DIR, "market_growth_chart_bw_1765140893793.png")

def create_presentation():
    prs = Presentation()

    # Helper to add a slide with title and content
    def add_slide(layout_index, title_text, content_text=None):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        title = slide.shapes.title
        title.text = title_text
        
        if content_text:
            # Try to find the content placeholder
            if len(slide.placeholders) > 1:
                content = slide.placeholders[1]
                content.text = content_text
            else:
                # Fallback if no content placeholder
                txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(5))
                tf = txBox.text_frame
                tf.text = content_text
        return slide

    # 1. Title Slide
    slide_layout = prs.slide_layouts[0] # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = "Track-subs"
    subtitle.text = "Revolutionizing Subscription Management in India\nInvestor Pitch Deck"
    
    # Add Hero Image to Title Slide
    if os.path.exists(IMG_TITLE_HERO):
        left = Inches(0)
        top = Inches(0)
        pic = slide.shapes.add_picture(IMG_TITLE_HERO, left, top, width=prs.slide_width, height=prs.slide_height)
        # Move picture to back so text is visible (this is a hack, pptx doesn't support z-order easily, 
        # so we re-add text boxes on top if needed, or just place image carefully. 
        # For now, let's just place it as a background and assume text might be covered, 
        # or better, place it on the side or bottom)
        # Actually, let's make it a side image for better visibility of text
        pic.left = Inches(0)
        pic.top = Inches(2.5)
        pic.height = Inches(5)
        pic.width = prs.slide_width

    # 2. The Problem
    slide = add_slide(1, "The Problem: Subscription Chaos")
    content = slide.placeholders[1]
    content.text = "• Average Indian user has 5-10 active subscriptions (OTT, Food, Utilities).\n" \
                   "• 'Free Trials' often turn into unwanted monthly charges.\n" \
                   "• Hidden UPI autopay mandates are hard to track.\n" \
                   "• Result: Wasted money and financial anxiety."
    
    if os.path.exists(IMG_PROBLEM):
        slide.shapes.add_picture(IMG_PROBLEM, Inches(5), Inches(2), height=Inches(4))

    # 3. The Solution: Track-subs
    slide = add_slide(1, "The Solution: Intelligent Tracking")
    content = slide.placeholders[1]
    content.text = "• Automated Detection: Scans bank statements for recurring patterns.\n" \
                   "• Smart Alerts: Notify before renewal dates.\n" \
                   "• One-Click Cancel: Actionable insights to stop leakage.\n" \
                   "• India-First: Optimized for UPI, RuPay, and local merchants."
    
    if os.path.exists(IMG_SOLUTION):
        slide.shapes.add_picture(IMG_SOLUTION, Inches(5), Inches(2), height=Inches(4))

    # 4. Market Opportunity
    slide = add_slide(1, "Market Opportunity")
    content = slide.placeholders[1]
    content.text = "• India's Subscription Economy is booming (CAGR 20%+).\n" \
                   "• 500M+ Digital Payment Users.\n" \
                   "• Rising adoption of SaaS and OTT platforms.\n" \
                   "• No dominant player in the personal subscription management space."
    
    if os.path.exists(IMG_MARKET):
        slide.shapes.add_picture(IMG_MARKET, Inches(5), Inches(2), height=Inches(4))

    # 5. Product Demo (Mention)
    slide = add_slide(1, "Product Status")
    content = slide.placeholders[1]
    content.text = "• Functional MVP Live.\n" \
                   "• Backend: FastAPI with ML-based detection logic.\n" \
                   "• Frontend: React + Tailwind CSS.\n" \
                   "• Security: OTP-based authentication implemented.\n" \
                   "• Ready for beta testing."

    # 6. Business Model
    slide = add_slide(1, "Business Model")
    content = slide.placeholders[1]
    content.text = "• Freemium: Basic tracking is free.\n" \
                   "• Premium: Advanced analytics, concierge cancellation service.\n" \
                   "• Affiliate: Commission from recommending better/cheaper plans.\n" \
                   "• Data Insights: Aggregated trends for merchants (B2B)."

    # 7. The Team
    slide = add_slide(1, "The Team")
    content = slide.placeholders[1]
    content.text = "• Himanshu Singhal: Founder & Lead Developer.\n" \
                   "• Passionate about Fintech and solving real-world problems.\n" \
                   "• Experienced in Full Stack Development and AI."

    # Save
    prs.save(OUTPUT_FILE)
    print(f"Presentation saved to {OUTPUT_FILE}")

if __name__ == "__main__":
    create_presentation()
